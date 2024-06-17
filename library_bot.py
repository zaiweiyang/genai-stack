import os
import json
import grpc
import file_service_pb2
import file_service_pb2_grpc
import requests
import datetime
import threading
import concurrent.futures
from cassandra.cluster import Cluster
from cassandra.auth import PlainTextAuthProvider

import streamlit as st
from langchain.chains import RetrievalQA 
from PyPDF2 import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation
import io
from langchain.callbacks.base import BaseCallbackHandler
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import Neo4jVector
from langchain.schema.document import Document
from langchain_core.runnables import RunnableConfig

from streamlit.logger import get_logger
from chains import (
    load_embedding_model,
    load_llm,
)

import phoenix as px
from phoenix.trace.langchain import LangChainInstrumentor
import logging
logging.basicConfig(level=logging.INFO)

import openai

# session = px.launch_app()
# set_global_handler("arize_phoenix")

# load api key lib
from dotenv import load_dotenv
# from custom_widgets.layout import Accordion

# from langchain.vectorstores import Chroma
from langchain_community.vectorstores import Chroma

import uuid
NAMESPACE_UUID = uuid.uuid4()

# Operation status flag
operation_in_progress = False

load_dotenv(".env")

url = os.getenv("NEO4J_URI")
username = os.getenv("NEO4J_USERNAME")
password = os.getenv("NEO4J_PASSWORD")
ollama_base_url = os.getenv("OLLAMA_BASE_URL")
embedding_model_name = os.getenv("EMBEDDING_MODEL")
llm_name = os.getenv("LLM")
backup_api_svc_url = os.getenv("BACKUP_API_SVC_URL")
phoenix_tracing_endpoint = os.getenv("PHOENIX_COLLECTOR_ENDPOINT")
grpc_server_url = os.getenv("GRPC_SERVER_URL", "host.docker.internal:50051")

# Cassandra connection details
cassandra_host = os.getenv("CASSANDRA_HOST")
cassandra_username = os.getenv("CASSANDRA_USERNAME")
cassandra_password = os.getenv("CASSANDRA_PASSWORD")

# # Connect to Cassandra
# logging.info(f"cassandra_host: {cassandra_host}")

auth_provider = PlainTextAuthProvider(username=cassandra_username, password=cassandra_password)
cluster = Cluster([cassandra_host], auth_provider=auth_provider)
session = cluster.connect()

# Create a keyspace and table if they don't exist
session.execute("""
    CREATE KEYSPACE IF NOT EXISTS documents
    WITH REPLICATION = { 'class' : 'SimpleStrategy', 'replication_factor' : 1 }
""")
session.execute("""
    CREATE TABLE IF NOT EXISTS documents.summaries (
        id UUID PRIMARY KEY,
        title TEXT,
        abstract TEXT,
        summary TEXT,
        uploaded_to TEXT,
        file_type TEXT,
        file_name TEXT
    )
""")

# Remapping for Langchain Neo4j integration
os.environ["NEO4J_URL"] = url

#### PHOENIX Endpoint should be configured via environment variables `PHOENIX_HOST`, `PHOENIX_PORT`, or `PHOENIX_COLLECTOR_ENDPOINT`.
# os.environ["PHOENIX_PORT"] = "6006"
# os.environ["PHOENIX_HOST"] = "192.168.31.32"
os.environ["PHOENIX_COLLECTOR_ENDPOINT"] = phoenix_tracing_endpoint
LangChainInstrumentor().instrument()

logger = get_logger(__name__)

embeddings, dimension = load_embedding_model(
    embedding_model_name, config={"ollama_base_url": ollama_base_url}, logger=logger
)

class Accordion:
    def __init__(self, label, expanded=True):
        self.label = label
        self.content = ""
        self.expanded = expanded
        self.container = st.empty()

    def markdown(self, text):
        self.content = text
        self.update()

    def update(self):
        with self.container.expander(self.label, expanded=self.expanded):
            st.markdown(self.content)

class StreamHandler(BaseCallbackHandler):
    def __init__(self, container, name):
        self.container = container
        self.name = name
        self.text = ""
        self.lock = threading.Lock()
        self.llm_accordion = Accordion(label=f"**{name} LLM Answer**")
        self.retriever_accordion = Accordion(label=f"**{name} Retriever Output**", expanded=False)
        self.retriever_text = ""
        self.llm_accordion.update()
        self.retriever_accordion.update()
        self.chain_text = ""

    def on_llm_new_token(self, token: str, **kwargs) -> None:
        with self.lock:
            self.text += token
            self.llm_accordion.markdown(self.text)

    def on_retriever_end(self, documents, **kwargs):
        with self.lock:
            document_detail = ""
            for idx, doc in enumerate(documents):
                self.retriever_text += f"\n\n**Results from Document[{idx}]:**\n\n"
                self.retriever_text += doc.page_content
                document_detail += json.dumps(doc.metadata, indent=2)
                self.retriever_accordion.markdown(self.retriever_text)

    def on_chain_end(self, outputs, **kwargs):
        with self.lock:
            if self.retriever_text == "":
                self.chain_text += json.dumps(outputs, indent=2)
                self.chain_text += f"\n\n"
                self.retriever_accordion.markdown(self.chain_text)

llm = load_llm(llm_name, logger=logger, config={"ollama_base_url": ollama_base_url})

def escape_quotes(text):
    if text is not None:
        return text.replace("'", "''")
    return text

def summarize_text_with_openai(text):
    client = openai.OpenAI(api_key=os.getenv("OPENAI_API_KEY"))
    try:
        openai_model = os.getenv("CHATGPT_MODEL", "gpt-4")

        logger.info(f"OpenAI Model: {openai_model}")
        response = client.chat.completions.create(
            model=openai_model,
            messages=[
                {"role": "system", "content": "You are a helpful assistant."},
                {"role": "user", "content": f"Summarize the following text:\n\n{text}"}
            ],
            max_tokens=2300,
            stop=None,
            temperature=0.5,
            top_p=0.9
        )
        summary = response.choices[0].message.content.strip()
        return summary
    except Exception as e:
        logger.error(f"Failed to summarize text with OpenAI: {str(e)}")
        return "None"

# Keep the latest record only and remove other duplicated records with the same file_name from the database.
def store_metadata_in_cassandra(doc_id, title, abstract, summary, uploaded_to, file_type, file_name):
    try:
        title = escape_quotes(title)
        abstract = escape_quotes(abstract)
        summary = escape_quotes(summary)
        uploaded_to = escape_quotes(uploaded_to)
        file_type = escape_quotes(file_type)
        file_name = escape_quotes(file_name)

        # Select existing records with the same file_name
        rows = session.execute("""
            SELECT id FROM documents.summaries WHERE file_name = %s ALLOW FILTERING
        """, (file_name,))

        # Delete each selected record using the full primary key
        for row in rows:
            session.execute("""
                DELETE FROM documents.summaries WHERE id = %s
            """, (row.id,))

        # Insert the new record
        session.execute("""
            INSERT INTO documents.summaries (id, title, abstract, summary, uploaded_to, file_type, file_name)
            VALUES (%s, %s, %s, %s, %s, %s, %s)
        """, (doc_id, title, abstract, summary, uploaded_to, file_type, file_name))

    except Exception as e:
        logger.error(f"Store metadata in cassandra error: {e}")
        return []

def download_file(file_path):
    channel = grpc.insecure_channel(grpc_server_url)
    stub = file_service_pb2_grpc.FileServiceStub(channel)

    response = stub.GetFile(file_service_pb2.FileRequest(file_path=file_path))
    return response.content

def handle_file_upload(file_content, file_name, save_to_chroma, save_to_neo4j, file_type):
    if file_type == 'pdf':
        file_like_object = io.BytesIO(file_content)
        pdf_reader = PdfReader(file_like_object)
        title = pdf_reader.metadata.get('/Title', "No Title Available") if pdf_reader.metadata else "No Metadata Found"
        text = "".join((page.extract_text() or "") for page in pdf_reader.pages)
    elif file_type == 'docx':
        file_like_object = io.BytesIO(file_content)
        docx_document = DocxDocument(file_like_object)
        text = "\n\n".join([paragraph.text for paragraph in docx_document.paragraphs])
        title = file_name
    elif file_type == 'pptx':
        file_like_object = io.BytesIO(file_content)
        pptx_document = Presentation(file_like_object)
        text = "\n\n".join([shape.text for slide in pptx_document.slides for shape in slide.shapes if hasattr(shape, "text")])
        title = file_name
    elif file_type == 'txt':
        text = file_content.decode("utf-8")
        title = file_name

    abstract = text[:500]  # First 500 characters as an abstract
    summary = summarize_text_with_openai(text)  # Generate summary with OpenAI
    logger.info(f"<<<Got document summary <{summary}>.")

    text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200, length_function=len)
    chunks = text_splitter.split_text(text=text)
    logger.info(f"{len(chunks)} chunks are generated from the {file_type} <{file_name}>.")

    ids = [str(uuid.uuid5(NAMESPACE_UUID, chunk)) for chunk in chunks]
    unique_ids = list(set(ids))
    seen_ids = set()
    unique_chunks = [chunk for chunk, id in zip(chunks, ids) if id not in seen_ids and (seen_ids.add(id) or True)]

    uploaded_to = []
    if save_to_chroma:
        logger.info(f"{len(unique_chunks)} unique chunks saved to Chroma")
        chroma_db = Chroma.from_texts(texts=unique_chunks, embedding=embeddings, ids=unique_ids, persist_directory="/data", collection_name=f"{file_type}_bot")
        chroma_db.persist()
        uploaded_to.append('Chroma')

    if save_to_neo4j:
        logger.info(f"{len(unique_chunks)} unique chunks saved to Neo4j")
        vectorstore = Neo4jVector.from_texts(
            chunks,
            url=url,
            username=username,
            password=password,
            embedding=embeddings,
            ids=unique_ids,
            index_name=f"{file_type}_bot",
            node_label=f"{file_type.capitalize()}BotChunk"
        )
        uploaded_to.append('Neo4j')

    store_metadata_in_cassandra(uuid.uuid4(), title, abstract, summary, ','.join(uploaded_to), file_type, file_name)

def scan_directories(directories, save_to_chroma, save_to_neo4j):
    channel = grpc.insecure_channel(grpc_server_url)
    stub = file_service_pb2_grpc.FileServiceStub(channel)

    for directory in directories:
        response = stub.ListFiles(file_service_pb2.DirectoryRequest(directory_path=directory))
        for file_name in response.files:
            file_path = os.path.join(directory, file_name)
            file_content = download_file(file_path)
            file_type = os.path.splitext(file_name)[1][1:]  # Get the file extension
            handle_file_upload(file_content, file_name, save_to_chroma, save_to_neo4j, file_type)

def load_config():
    config_path = './config.json'
    with open(config_path, 'r') as file:
        config = json.load(file)
    
    return config.get('directories', [])

def display_uploaded_pdfs():
    st.subheader("Previously uploaded documents:")
    rows = session.execute("SELECT file_name, title, abstract, uploaded_to FROM documents.summaries")
    for row in rows:
        uploaded_to = ', '.join(set(row.uploaded_to.split(',')))
        st.markdown(f"**{row.file_name}** - Title: {row.title}, Abstract: {row.abstract[:150]}... (Uploaded to: {uploaded_to})")

def run_qa_section():
    st.header("Document Query Assistant")
    
    col1, col2 = st.columns(2)

    chroma_db = Chroma(
        persist_directory="/data", 
        collection_name="pdf_bot",
        embedding_function=embeddings
    )
    chroma_retriever = chroma_db.as_retriever(search_kwargs={"k": 6})

    qa_chroma = RetrievalQA.from_chain_type(
        llm=llm, chain_type="stuff", retriever=chroma_retriever, return_source_documents=True
    )
    
    neo4j_retriever = Neo4jVector(
        url=url,
        username=username,
        password=password,
        embedding=embeddings,
        index_name="pdf_bot",
        node_label="PdfBotChunk"
    ).as_retriever()
    qa_neo4j = RetrievalQA.from_chain_type(
        llm=llm, chain_type="stuff", retriever=neo4j_retriever
    )

    db_col1, db_col2 = st.columns(2)
    qa_neo4j_selected = True
    qa_chroma_selected = True
    with db_col1:
        qa_neo4j_selected = st.checkbox("Query against Neo4j Graph Database", value=True)
    with db_col2:
        qa_chroma_selected = st.checkbox("Query against Chroma Vector Database", value=True)

    query = st.text_input("Ask a question:")
    if query:
        query_dict = {"query": query}
        if qa_neo4j_selected and qa_chroma_selected: 
            col1, col2 = st.columns(2)
            with col1:
                st.markdown("### Neo4j Database Response")
                stream_handler_neo4j = StreamHandler(st.empty(), "Neo4j")
                config = RunnableConfig({"callbacks": [stream_handler_neo4j],"run_name":"Neo4j"})
                result = qa_neo4j.invoke(query_dict, config=config)
            with col2:
                st.markdown("### Chroma Database Response")
                stream_handler_chroma = StreamHandler(st.empty(), "Chroma")
                config = RunnableConfig({"callbacks": [stream_handler_chroma],"run_name":"Chroma"})
                result = qa_chroma.invoke(query_dict, config=config)

        elif qa_neo4j_selected:
            st.markdown("### Neo4j Database Response")
            stream_handler_neo4j = StreamHandler(st.empty(), "Neo4j")
            config = RunnableConfig({"callbacks": [stream_handler_neo4j],"run_name":"Neo4j"})
            result = qa_neo4j.invoke(query_dict, config=config)
        elif qa_chroma_selected:
            st.markdown("### Chroma Database Response")
            stream_handler_chroma = StreamHandler(st.empty(), "Chroma")
            config = RunnableConfig({"callbacks": [stream_handler_chroma],"run_name":"Chroma"})
            result = qa_chroma.invoke(query_dict, config=config)
        else:
            st.markdown("### No Database is selected for the query")

def get_backup_tags():
    try:
        response = requests.get(f'{backup_api_svc_url}/backups')
        if response.status_code == 200:
            backup_tags = response.json()
            return backup_tags
        else:
            logger.error(f"Failed to retrieve backup tags: {response.status_code} {response.text}")
            return []
    except Exception as e:
        logger.error(f"Error retrieving backup tags: {e}")
        return []

def manage_backups():
    global operation_in_progress
    st.header("Database Management")
    col1, col2 = st.columns(2)
    backup_tags = get_backup_tags()
    
    with col1:
        tag = st.text_input('Enter a tag for this backup(optional):', key='backup_tag', disabled=operation_in_progress)
        if st.button('Backup Database', key='backup_button', disabled=operation_in_progress):
            if not tag:
                tag = datetime.datetime.now().strftime("%Y%m%d-%H%M%S")
            if tag in backup_tags:
                st.error("This tag already exists. Please use a different tag.")
            else:
                response = backup_database(tag)
                message = response.get('message', 'Backup failed due to an unexpected error.')
                st.text(message)
                
    with col2:
        tag = st.selectbox('Choose a tag to restore from:', backup_tags, key='restore_tag', disabled=operation_in_progress)
        if st.button('Restore Database', key='restore_button', disabled=operation_in_progress):
            response = restore_database(tag)
            message = response.get('message', 'Restore failed due to an unexpected error.')
            st.text(message)

def backup_database(tag=None):
    global operation_in_progress
    operation_in_progress = True
    data = {"tag": tag} if tag else {}
    response = requests.post(f'{backup_api_svc_url}/backup', json=data)

    operation_in_progress = False
    if response.status_code == 200:
        return response.json()
    else:
        return {"message": f"Backup failed with status {response.status_code}: {response.text}"}

def restore_database(tag):
    global operation_in_progress
    operation_in_progress = True
    data = {"tag": tag}
    response = requests.post(f'{backup_api_svc_url}/restore', json=data)
    operation_in_progress = False
    if response.status_code == 200:
        return response.json()
    else:
        return {"message": f"Restore failed with status {response.status_code}: {response.text}"}

def main():
    st.set_page_config(layout="wide")

    st.header("📄Demo: Generative AI enriched with local data. \nLocally hosted LLM Model and local Knowledge Base generated from documents.")

    directories = load_config()
    logging.info(f"loaded directories: {directories}")

    col1, col2 = st.columns([3, 1])

    with col1:
        save_to_chroma = True
        save_to_neo4j = True
        db_col1, db_col2 = st.columns([1, 1])  
        with db_col1:
            save_to_neo4j = st.checkbox("Save to Neo4j Graph Database", value=True)
        with db_col2:
            save_to_chroma = st.checkbox("Save to Chroma Vector Database", value=True)
        
        if not operation_in_progress:
            pdf = st.file_uploader("Upload your PDF", type="pdf")
            if pdf:
                handle_file_upload(pdf.read(), pdf.name, save_to_chroma, save_to_neo4j, 'pdf')

            docx = st.file_uploader("Upload your DOCX", type="docx")
            if docx:
                handle_file_upload(docx.read(), docx.name, save_to_chroma, save_to_neo4j, 'docx')

            pptx = st.file_uploader("Upload your PPTX", type="pptx")
            if pptx:
                handle_file_upload(pptx.read(), pptx.name, save_to_chroma, save_to_neo4j, 'pptx')

            txt = st.file_uploader("Upload your TXT", type="txt")
            if txt:
                handle_file_upload(txt.read(), txt.name, save_to_chroma, save_to_neo4j, 'txt')

        if not operation_in_progress:
            run_qa_section() 
        
        st.markdown(f"### [Tracing the execution]({phoenix_tracing_endpoint})")

        manage_backups()

    with col2:
        display_uploaded_pdfs()

    scan_directories(directories, save_to_chroma, save_to_neo4j)

if __name__ == "__main__":
    main()
